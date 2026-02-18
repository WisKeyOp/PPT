"""
Quick verification script to check the generated presentation structure.
"""
from pptx import Presentation

# Load the generated presentation
prs = Presentation("data/outputs/EY.pptx")

print(f"\n=== Generated Presentation Analysis ===")
print(f"Total slides: {len(prs.slides)}\n")

for i, slide in enumerate(prs.slides, 1):
    layout_name = slide.slide_layout.name
    print(f"Slide {i}: {layout_name}")
    
    # Show first few placeholders with text
    for shape in slide.shapes:
        if hasattr(shape, "text_frame") and shape.text_frame.text.strip():
            text_preview = shape.text_frame.text[:80] + "..." if len(shape.text_frame.text) > 80 else shape.text_frame.text
            print(f"  - {text_preview}")
    print()

print("=== Layout Diversity Check ===")
layout_usage = {}
for slide in prs.slides:
    layout_name = slide.slide_layout.name
    layout_usage[layout_name] = layout_usage.get(layout_name, 0) + 1

for layout, count in sorted(layout_usage.items(), key=lambda x: x[1], reverse=True):
    print(f"  {layout}: {count} slide(s)")
