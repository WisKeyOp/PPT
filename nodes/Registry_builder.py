import json
import os
import hashlib
import sys
import argparse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

# --- CONFIGURATION & SEMANTICS ---
DUMMY_TEXT_KEYWORDS = ["lorem ipsum", "click to add", "insert text", "sample data", "text here"]

def get_file_hash(filepath):
    """Generates an MD5 hash to detect template changes."""
    hash_md5 = hashlib.md5()
    with open(filepath, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_md5.update(chunk)
    return hash_md5.hexdigest()

def clean_text(text):
    if not text: return ""
    # Standardizing characters for the LLM
    text = text.replace('Æ', "'").replace('á', ' ').replace('\u2013', '-').replace('\u2019', "'")
    return " ".join(text.split())

def get_semantic_role(shape):
    """Maps PPTX placeholder types to Pipeline 2 instructions."""
    if not shape.is_placeholder:
        if shape.has_table: return "STATIC_TABLE"
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE: return "STATIC_IMAGE"
        return "AD_HOC_TEXT_BOX"
    
    ph_type = shape.placeholder_format.type
    mapping = {
        PP_PLACEHOLDER.TITLE: "SLIDE_TITLE",
        PP_PLACEHOLDER.CENTER_TITLE: "MAIN_HEADING",
        PP_PLACEHOLDER.SUBTITLE: "SUB_HEADING",
        PP_PLACEHOLDER.BODY: "CONTENT_BODY",
        PP_PLACEHOLDER.PICTURE: "MEDIA_SLOT",
        PP_PLACEHOLDER.CHART: "DATA_VIS_SLOT",
        PP_PLACEHOLDER.TABLE: "TABLE_SLOT",
        PP_PLACEHOLDER.OBJECT: "GENERAL_OBJECT_SLOT"
    }
    return mapping.get(ph_type, "OTHER_PLACEHOLDER")

def is_concentric(inner_geo, outer_geo):
    """Checks if inner_geo is physically inside outer_geo (with small margin)."""
    # rect format: (left, top, width, height)
    ix1, iy1 = inner_geo[0], inner_geo[1]
    ix2, iy2 = ix1 + inner_geo[2], iy1 + inner_geo[3]
    
    ox1, oy1 = outer_geo[0], outer_geo[1]
    ox2, oy2 = ox1 + outer_geo[2], oy1 + outer_geo[3]
    
    # Check if inner box boundaries are within outer box boundaries
    return (ix1 >= ox1 - 1000 and iy1 >= oy1 - 1000 and 
            ix2 <= ox2 + 1000 and iy2 <= oy2 + 1000)

def extract_superior_registry(pptx_path):
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        return {"error": str(e)}

    sw, sh = prs.slide_width, prs.slide_height
    
    # 1. Map Sections for global context
    section_map = {}
    try:
        for section in prs.section_properties.sections:
            for s_id in section.slide_id_list:
                section_map[s_id] = section.name
    except: pass

    blueprint = {
        "registry_metadata": {
            "template_file": os.path.basename(pptx_path),
            "file_hash": get_file_hash(pptx_path),
            "dimensions_emu": {"width": sw, "height": sh},
            "aspect_ratio": round(sw/sh, 2)
        },
        "slides": []
    }

    for i, slide in enumerate(prs.slides, 1):
        slide_section = section_map.get(slide.slide_id, "Main Content")
        slide_data = {
            "slide_index": i,
            "section": slide_section,
            "layout": slide.slide_layout.name,
            "elements": []
        }

        # Process all shapes first to get geometry
        raw_elements = []
        for shape in slide.shapes:
            text = clean_text(shape.text) if hasattr(shape, "text") else ""
            is_dummy = any(kw in text.lower() for kw in DUMMY_TEXT_KEYWORDS)

            # Guard: any spatial property can be None for groups/connectors/decorators
            s_left   = shape.left   or 0
            s_top    = shape.top    or 0
            s_width  = shape.width  or 0
            s_height = shape.height or 0

            el = {
                "id": shape.shape_id,
                "role": get_semantic_role(shape),
                "name": shape.name,
                "text_content": text if not is_dummy else "",
                "format_hint": text if is_dummy else "",
                "geometry": {
                    "left": s_left, "top": s_top,
                    "width": s_width, "height": s_height,
                    "rel_x": round(s_left / sw, 4),
                    "rel_y": round(s_top / sh, 4),
                    "area_pct": round(((s_width * s_height) / (sw * sh)) * 100, 2)
                },
                "children_ids": []  # For concentric nesting
            }
            raw_elements.append(el)

        # 2. Resolve Hierarchy (Parent-Child concentric boxes)
        # Sort by area descending so we check smaller boxes against larger ones
        sorted_by_area = sorted(raw_elements, key=lambda x: x["geometry"]["area_pct"], reverse=True)
        final_elements = []
        nested_ids = set()

        for idx, outer in enumerate(sorted_by_area):
            if outer["id"] in nested_ids: continue
            
            o_geo = (outer["geometry"]["left"], outer["geometry"]["top"], 
                     outer["geometry"]["width"], outer["geometry"]["height"])
            
            for inner in sorted_by_area[idx+1:]:
                if inner["id"] in nested_ids: continue
                
                i_geo = (inner["geometry"]["left"], inner["geometry"]["top"], 
                         inner["geometry"]["width"], inner["geometry"]["height"])
                
                if is_concentric(i_geo, o_geo):
                    outer["children_ids"].append(inner["id"])
                    nested_ids.add(inner["id"])

        # Final filtering: remove zero-area artifacts
        slide_data["elements"] = [e for e in sorted_by_area if e["geometry"]["area_pct"] > 0.01]
        blueprint["slides"].append(slide_data)

    return blueprint

if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("template", help="Path to .pptx template")
    args = parser.parse_args()

    # --- THE REGISTRY MANAGER LOGIC ---
    registry_file = f"{os.path.splitext(args.template)[0]}_registry.json"
    current_hash = get_file_hash(args.template)
    
    needs_rebuild = True
    if os.path.exists(registry_file):
        with open(registry_file, 'r') as f:
            existing_data = json.load(f)
            if existing_data.get("registry_metadata", {}).get("file_hash") == current_hash:
                needs_rebuild = False

    if needs_rebuild:
        sys.stderr.write("🏗 Template changed or missing. Rebuilding Registry...\n")
        full_registry = extract_superior_registry(args.template)
        with open(registry_file, 'w', encoding='utf-8') as f:
            json.dump(full_registry, f, indent=4, ensure_ascii=False)
        print(json.dumps(full_registry, indent=2))
    else:
        sys.stderr.write("⚡ Template unchanged. Loading cached Registry.\n")
        with open(registry_file, 'r') as f:
            print(f.read())