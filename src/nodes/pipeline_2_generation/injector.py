# src/nodes/pipeline_2_generation/injector.py
"""
Pipeline 2 Node: Injector
Responsibility: Render exactly what it is told (deterministic PPTX rendering)
No markdown parsing, no heuristics, no decisions - pure XML rendering

Now includes:
- Background gradient/overlay rendering (visual, not just specs)
- Placeholder-specific header text filtering
"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from typing import Optional, Any, Dict, List
from pptx.dml.color import RGBColor
from src.core.state import PPTState, BackgroundImageSpec
from src.utils.ppt_helper import find_placeholder_by_id
from src.config import ENABLE_AUTOFIT_ROLES
import os
import re

# Patterns to detect and skip repeated template header text
HEADER_PATTERNS = [
    r"Presentation\s+title\s+Page\s+\d+",  # "Presentation title Page X"
    r"^\s*Page\s+\d+\s+\d{1,2}\s+\w+\s+\d{4}\s*$",  # "Page X DD Month YYYY"
    r"^\s*\d{1,2}\s+\w+\s+\d{4}\s*$"  # "DD Month YYYY" alone
]

# Role-based visual theming
ROLE_BACKGROUND_COLORS = {
    "TITLE": RGBColor(0, 51, 102),      # Deep blue
    "AGENDA": RGBColor(45, 45, 45),     # Dark gray
    "CLOSING": RGBColor(0, 102, 51),    # Deep green
}


def _is_template_header(text: str, placeholder_type: Optional[str] = None) -> bool:
    """
    Check if text matches template header patterns that should be skipped.
    
    Args:
        text: Text content to check
        placeholder_type: Type of placeholder (if available) for smarter filtering
        
    Returns:
        True if text matches a header pattern
    """
    text_stripped = text.strip() if text else ""
    if not text_stripped:
        return False
    
    # If we know it's a title/subtitle placeholder on slide 1, don't filter
    if placeholder_type in ["TITLE", "CENTER_TITLE", "SUBTITLE"]:
        return False
    
    for pattern in HEADER_PATTERNS:
        if re.search(pattern, text_stripped, re.IGNORECASE):
            return True
    
    return False


def _add_background_gradient(slide, background_spec: Dict[str, Any], slide_role: str) -> None:
    """
    Add a visual background to a slide with gradient and overlay.
    
    Since we don't have actual image retrieval yet, this creates:
    1. A full-slide gradient background based on role
    2. An overlay with appropriate opacity for text readability
    3. Stores image spec in notes for future image integration
    
    Args:
        slide: The slide object
        background_spec: Background image specification dict
        slide_role: The role of the slide (TITLE, AGENDA, etc.)
    """
    if not background_spec.get("enabled", False):
        return
    
    # Get slide dimensions
    prs = slide.part.package.presentation_part.presentation
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # Add full-bleed background rectangle
    left = top = 0
    
    # Get role-based color
    bg_color = ROLE_BACKGROUND_COLORS.get(slide_role, RGBColor(30, 30, 30))
    
    # Add background shape at the back
    from pptx.enum.shapes import MSO_SHAPE
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, slide_width, slide_height
    )
    
    # Apply gradient fill
    fill = background.fill
    fill.gradient()
    fill.gradient_angle = 45.0  # Diagonal gradient
    
    # Set gradient stops (darker at top-left, lighter at bottom-right)
    # Calculate lighter color (RGBColor doesn't have .red/.green/.blue attributes directly)
    # We'll use tuple unpacking from the base color
    r, g, b = bg_color
    fill.gradient_stops[0].color.rgb = bg_color
    fill.gradient_stops[1].color.rgb = RGBColor(
        min(255, r + 40),
        min(255, g + 40),
        min(255, b + 40)
    )
    
    # Remove line
    background.line.fill.background()
    
    # Send to back (z-order)
    # This ensures content appears on top
    slide.shapes._spTree.remove(background._element)
    slide.shapes._spTree.insert(2, background._element)  # Insert after background shapes
    
    # Add semi-transparent overlay for text readability
    overlay_opacity = background_spec.get("overlay_opacity", 0.4)
    
    overlay = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left, top, slide_width, slide_height
    )
    
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(20, 20, 20)  # Dark overlay
    overlay.fill.transparency = overlay_opacity  # Use spec opacity
    overlay.line.fill.background()
    
    # Send overlay to back (but above gradient)
    slide.shapes._spTree.remove(overlay._element)
    slide.shapes._spTree.insert(3, overlay._element)
    
    # Add image spec to slide notes for future integration
    keywords = background_spec.get("keywords", "")
    mood = background_spec.get("mood", "")
    composition = background_spec.get("composition", "")
    
    notes_text = f"""Background Image Specification:
- Keywords: {keywords}
- Mood: {mood}
- Composition: {composition}
- Overlay Opacity: {overlay_opacity:.0%}

Current: Role-based gradient background ({slide_role})
Future: Replace with actual image using keywords above
"""
    
    # Set notes
    if hasattr(slide, 'notes_slide'):
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = notes_text.strip()
    
    print(f"  ✓ Background applied: {slide_role} gradient with {overlay_opacity:.0%} overlay")


def find_placeholder_by_semantic_role(slide, semantic_role: str, slot_metadata: List[Dict]) -> Optional[Any]:
    """
    Find placeholder using semantic role with 3-tier fallback strategy.
    
    This function maps semantic content (title, bullets, body, image_query) to actual
    PowerPoint placeholders using stable identity from the registry.
    
    Fallback strategy:
    1. Try placeholder_idx (PowerPoint's canonical ID)
    2. Try placeholder_type (TITLE, BODY, PICTURE, etc.)
    3. Try placeholder_name ("Title 1", "Content Placeholder 2", etc.)
    
    Args:
        slide: The slide object to search
        semantic_role: The semantic role to find (title, bullets, body, image_query, footer)
        slot_metadata: List of slot dicts with stable identity fields
        
    Returns:
        Placeholder shape or None if not found
    """
    # Find slot(s) with matching semantic role
    matching_slots = [s for s in slot_metadata if s.get('semantic_role') == semantic_role]
    
    if not matching_slots:
        return None
    
    # Use the first matching slot (most layouts have one of each semantic role)
    target_slot = matching_slots[0]
    
    # Tier 1: Try placeholder_idx (most reliable)
    placeholder_idx = target_slot.get('placeholder_idx')
    if placeholder_idx is not None:
        try:
            for shape in slide.placeholders:
                if hasattr(shape, 'placeholder_format'):
                    if shape.placeholder_format.idx == placeholder_idx:
                        return shape
        except Exception:
            pass
    
    # Tier 2: Try placeholder_type (e.g., "TITLE (1)", "BODY (2)")
    placeholder_type = target_slot.get('placeholder_type', '')
    if placeholder_type:
        # Extract the type name (e.g., "TITLE" from "TITLE (1)")
        type_name = placeholder_type.split('(')[0].strip() if '(' in placeholder_type else placeholder_type
        try:
            for shape in slide.placeholders:
                if hasattr(shape, 'placeholder_format'):
                    shape_type_str = str(shape.placeholder_format.type)
                    if type_name in shape_type_str:
                        return shape
        except Exception:
            pass
    
    # Tier 3: Try placeholder_name ("Title 1", "Content Placeholder 2", etc.)
    placeholder_name = target_slot.get('placeholder_name')
    if placeholder_name:
        try:
            for shape in slide.placeholders:
                if hasattr(shape, 'name') and shape.name == placeholder_name:
                    return shape
        except Exception:
            pass
    
    return None


def surgical_injection_node(state: PPTState):
    """
    Pipeline 2 – Final Node
    Deterministic PPTX renderer with SEMANTIC MAPPING.
    
    Maps semantic content (title, bullets, body, image_query) to placeholders
    using stable identity (placeholder_idx, type, name) instead of arbitrary slot IDs.
    
    This prevents overlap bugs caused by LLM-generated slot IDs not matching
    PowerPoint's actual placeholder indices.
    
    Features:
    - Semantic-to-placeholder mapping with 3-tier fallback
    - Renders role-based background gradients for enabled slides
    - Filters template header text intelligently  
    - Applies white text on backgrounds for readability
    """
    primary_master_path = state["primary_master_path"]
    manifest = state["manifest"]
    
    prs = Presentation(primary_master_path)
    print(f"--- Injector: Rendering {len(manifest)} slides ---")
    
    for i, slide_def in enumerate(manifest):
        layout_idx = slide_def["layout_index"]
        semantic_content = slide_def["content"]  # Now contains semantic fields, not slot_id-based
        background_spec = slide_def.get("background_image", {})
        slide_role = slide_def.get("slide_role", "CONTENT")
        semantic_mapping = slide_def.get("_semantic_mapping", {})  # Slot metadata from Writer
        
        # Add slide with specified layout
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        
        # Apply background gradient FIRST (before content, so it's behind everything)
        has_background = background_spec.get("enabled", False)
        if has_background:
            _add_background_gradient(slide, background_spec, slide_role)
        
        # Track which semantic fields were successfully rendered
        rendered_fields = []
        
        # Check if content is semantic (new) or slot_id-based (legacy for Beautifier)
        # Semantic content NOW has styled runs from Beautifier (not raw text)
        is_semantic = slide_def.get("_is_semantic", False)
        
        if is_semantic:
            # NEW SEMANTIC MAPPING APPROACH (with styled runs from Beautifier)
            # Beautifier now provides styled runs in slot_id format
            
            # Log content count for debugging empty slides
            print(f"  Slide {i+1}: Processing {len(semantic_content)} semantic slots")
            
            for slot_id, style in semantic_content.items():
                semantic_role = style.get("semantic_role")
                
                # DEFENSIVE: Skip image_query if somehow present
                if semantic_role == 'image_query':
                    print(f"  ⊘ Skipping image_query placeholder (NO_IMAGE mode)")
                    continue
                
                try:
                    # Find placeholder using slot_id (Beautifier provides slot_id-based output)
                    shape = find_placeholder_by_id(slide, int(slot_id))
                    
                    if not shape or not hasattr(shape, "text_frame"):
                        print(f"  ⊘ Skipped slot {slot_id} on slide {i+1}: placeholder not found")
                        continue
                    
                    # Get text frame and existing content
                    tf = shape.text_frame
                    existing_text = (tf.text or "").strip()
                    
                    # Filter template headers and instruction blocks (copy pattern from legacy path)
                    placeholder_type = None
                    if hasattr(shape, "placeholder_format"):
                        try:
                            placeholder_type = shape.placeholder_format.type
                        except Exception:
                            pass
                    
                    if _is_template_header(existing_text, placeholder_type):
                        print(f"  ⊘ Filtered template header/instruction in slot {slot_id} on slide {i+1}")
                        # Proceed with clearing and injection
                    
                    # CRITICAL: Validate content exists BEFORE clearing (prevents blank slides)
                    runs = style.get("runs") or []
                    bullets = style.get("bullets") or []
                    
                    has_runs = any((r.get("text") or "").strip() for r in runs)
                    has_bullets = bool(bullets)
                    
                    if not (has_runs or has_bullets):
                        print(f"  ⊘ Skipped slot {slot_id} on slide {i+1}: no content to inject")
                        continue  # Preserve template placeholder, don't create blank slide
                    
                    # Now safe to clear - we have content to inject
                    tf.clear()
                    
                    # Apply vertical anchor if specified (title slides)
                    if style.get("vertical_anchor") == "MIDDLE":
                        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    
                    # Get alignment from Beautifier (None means respect Master Slide default)
                    alignment_str = style.get("alignment")
                    if alignment_str == "CENTER":
                        target_alignment = PP_ALIGN.CENTER
                    elif alignment_str == "RIGHT":
                        target_alignment = PP_ALIGN.RIGHT
                    elif alignment_str == "LEFT":
                        target_alignment = PP_ALIGN.LEFT
                    else:
                        target_alignment = None  # Preserve Master Slide default
                    
                    # Render content based on structure
                    if "bullets" in style:
                        # BULLETS: One paragraph per bullet item
                        bullet_items = style.get("bullets", [])
                        
                        for bullet_idx, bullet_item in enumerate(bullet_items):
                            if bullet_idx == 0:
                                # DEFENSIVE: Ensure paragraph exists after clear()
                                p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                            else:
                                p = tf.add_paragraph()
                            
                            p.level = 0  # Bullet level
                            
                            # Render styled runs for this bullet
                            for run_data in bullet_item.get("runs", []):
                                run_text = run_data.get("text", "")
                                if run_text:
                                    r = p.add_run()
                                    r.text = run_text
                                    r.font.bold = run_data.get("bold", False)
                                    r.font.italic = run_data.get("italic", False)
                                    
                                    # Only apply font size if explicitly provided (not None)
                                    font_size = style.get("font_size")
                                    if font_size is not None:
                                        r.font.size = Pt(font_size)
                                    
                                    if has_background:
                                        r.font.color.rgb = RGBColor(255, 255, 255)
                        
                        # Apply alignment only if explicitly set (prevents overriding Master defaults)
                        if target_alignment is not None:
                            for p in tf.paragraphs:
                                p.alignment = target_alignment
                    
                    elif "runs" in style:
                        # SINGLE TEXT FIELD (title, body, footer)
                        # DEFENSIVE: Ensure paragraph exists after clear()
                        p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
                        
                        # Render styled runs
                        for run_data in style.get("runs", []):
                            run_text = run_data.get("text", "")
                            if run_text:
                                r = p.add_run()
                                r.text = run_text
                                r.font.bold = run_data.get("bold", False)
                                r.font.italic = run_data.get("italic", False)
                                
                                # Only apply font size if explicitly provided (not None)
                                font_size = style.get("font_size")
                                if font_size is not None:
                                    r.font.size = Pt(font_size)
                                
                                if has_background:
                                    r.font.color.rgb = RGBColor(255, 255, 255)
                        
                        # Apply alignment only if explicitly set (prevents overriding Master defaults)
                        if target_alignment is not None:
                            p.alignment = target_alignment
                    
                    # Optional: Enable autofit for specific roles (caption, circular_text)
                    semantic_role = style.get("semantic_role", "")
                    if semantic_role in ENABLE_AUTOFIT_ROLES:
                        try:
                            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                        except Exception:
                            pass  # Silently fail if not supported
                    
                    rendered_fields.append(slot_id)
                    print(f"  ✓ Rendered slot {slot_id} ({semantic_role}) on slide {i+1}")
                    
                except Exception as e:
                    print(f"  ⚠️  Failed to render slot {slot_id} on slide {i+1}: {e}")
            
            # Warn if no slots were rendered (helps debug empty slides)
            if not rendered_fields:
                print(f"  ⚠️  Warning: No content rendered on slide {i+1} (slide may appear empty)")
        else:
            # LEGACY SLOT_ID-BASED APPROACH (for backward compatibility with Beautifier)
            for slot_id, style in semantic_content.items():
                try:
                    # Find placeholder by numeric ID
                    shape = find_placeholder_by_id(slide, int(slot_id))
                    
                    if not shape or not hasattr(shape, "text_frame"):
                        continue
                    
                    # Get text frame and clear existing content
                    tf = shape.text_frame  # type: ignore - hasattr check above ensures this exists
                    
                    # Check for template header text in existing placeholder
                    existing_text = tf.text if hasattr(tf, 'text') else ""
                    
                    # Determine placeholder type for smarter filtering
                    placeholder_type = None
                    if hasattr(shape, 'placeholder_format'):
                        try:
                            placeholder_type = str(shape.placeholder_format.type)
                        except:
                            pass
                    
                    if _is_template_header(existing_text, placeholder_type):
                        # Clear the template header
                        print(f"  ⊘ Filtered template header in slot {slot_id}")
                        tf.clear()
                        # Check if we have actual content to render
                        runs_to_render = style.get("runs", [])
                        if not runs_to_render or all(not r.get("text", "").strip() for r in runs_to_render):
                            # No content to render, leave it empty
                            continue
                    else:
                        # Normal case: clear and populate
                        tf.clear()
                    
                    # Get or create paragraph (ensure at least one exists after clear)
                    if not tf.paragraphs:
                        p = tf.add_paragraph()
                    else:
                        p = tf.paragraphs[0]
                    
                    # Set alignment from style spec (None means respect Master default)
                    alignment_str = style.get("alignment")
                    if alignment_str == "CENTER":
                        target_alignment = PP_ALIGN.CENTER
                    elif alignment_str == "RIGHT":
                        target_alignment = PP_ALIGN.RIGHT
                    elif alignment_str == "LEFT":
                        target_alignment = PP_ALIGN.LEFT
                    else:
                        target_alignment = None
                    
                    # Render each styled run
                    has_content = False
                    for run_data in style.get("runs", []):
                        run_text = run_data.get("text", "")
                        
                        # Skip if this run matches a header pattern
                        if _is_template_header(run_text, None):
                            continue
                        
                        if run_text:  # Only add non-empty runs
                            r = p.add_run()
                            r.text = run_text
                            r.font.bold = run_data.get("bold", False)
                            r.font.italic = run_data.get("italic", False)
                            
                            # Only apply font size if explicitly provided (not None)
                            font_size = style.get("font_size")
                            if font_size is not None:
                                r.font.size = Pt(font_size)
                            
                            # Use white text on backgrounds for readability
                            if has_background:
                                r.font.color.rgb = RGBColor(255, 255, 255)
                            
                            has_content = True
                    
                    # Apply alignment only if explicitly set
                    if target_alignment is not None:
                        p.alignment = target_alignment
                    
                    if has_content:
                        rendered_fields.append(slot_id)
                    
                except (ValueError, TypeError) as e:
                    print(f"  ⚠️  Warning: Failed to render slot {slot_id} on slide {i+1}: {e}")
        
        # Log rendering summary
        bg_status = "with gradient background" if has_background else "no background"
        render_type = "semantic" if is_semantic else "slot_id"
        print(f"  ✓ Slide {i + 1} ({slide_role}): {len(rendered_fields)} {render_type} fields filled, {bg_status}")
    
    # Save output - use the path from state if provided, otherwise use default
    output_path = state.get("final_file_path") or "data/outputs/final_deck.pptx"
    
    # Ensure the output directory exists
    output_dir = os.path.dirname(output_path)
    if output_dir:
        os.makedirs(output_dir, exist_ok=True)
    
    prs.save(output_path)
    
    print(f"--- Injector: Saved presentation to {output_path} ---")
    
    return {"final_file_path": output_path}