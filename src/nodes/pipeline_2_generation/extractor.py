"""
Pipeline 2 Node: Extractor
Responsibility: Ingest any input format → produce raw semantic content
Supports: .docx, .pptx, .txt, raw text
"""
import os
from docx import Document
from pptx import Presentation
from src.utils.auth_helper import get_llm
from src.core.state import PPTState


def extract_text_from_file(file_path: str) -> str:
    """
    Extracts raw text from supported file formats.
    
    Args:
        file_path: Path to .docx, .pptx, or .txt file
        
    Returns:
        Extracted text content
        
    Raises:
        ValueError: If file type is not supported
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == ".docx":
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    
    if ext == ".pptx":
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text.append(shape.text)
        return "\n".join(text)
    
    if ext == ".txt":
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    
    raise ValueError(f"Unsupported file type: {ext}. Supported formats: .docx, .pptx, .txt")


def extract_context_node(state: PPTState):
    """
    Pipeline 2 – Node 1
    Converts raw input (file path or text) into a structured content map.
    
    Detects if input is a file path and extracts text accordingly.
    Uses LLM to synthesize structured narrative covering:
    - Vision
    - Problem
    - Solution
    - Technical Details
    - Milestones
    """
    raw_input = state["raw_docs"]
    
    # Detect file path vs raw text
    if isinstance(raw_input, str) and os.path.exists(raw_input):
        print(f"--- Extractor: Reading file {os.path.basename(raw_input)} ---")
        text = extract_text_from_file(raw_input)
        print(f"--- Extractor: Extracted {len(text)} characters from file ---")
    else:
        text = raw_input
        print(f"--- Extractor: Processing raw text input ({len(text)} chars) ---")
    
    llm = get_llm(deployment_name="gpt-4", temperature=0)
    
    prompt = f"""
Analyze the following documentation and produce a structured content map.

Your analysis should cover these key areas:
1. **Vision**: Overarching goals and mission
2. **Problem**: Specific challenges and pain points being addressed
3. **Solution**: Proposed approach and key features
4. **Technical Details**: Architecture, specifications, implementation details
5. **Milestones**: Roadmap, timeline, next steps

DOCUMENTATION:
{text}

Provide a comprehensive, well-structured summary that captures the narrative arc 
from vision through technical execution. This will be used to generate presentation slides.
"""
    
    response = llm.invoke(prompt)
    print(f"--- Extractor: Generated content map ({len(response.content)} chars) ---")
    
    return {"content_map": response.content}
