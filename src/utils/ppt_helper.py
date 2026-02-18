# src/utils/ppt_helper.py
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE
import re
 
def find_placeholder_by_id(slide, idx):
    """
    Finds a specific placeholder on a given slide by its unique idx.
    """
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == idx:
            return shape
    return None


def map_shape_type(shape):
    """
    Map PowerPoint shape type to simplified category.
    
    Returns:
        - "oval" for circles and ellipses (MSO_SHAPE_TYPE.OVAL)
        - "rectangle" for rectangles (MSO_SHAPE_TYPE.PLACEHOLDER, AUTO_SHAPE, etc.)
        - "other" for complex/unsupported shapes
    """
    try:
        if not hasattr(shape, 'shape_type'):
            return "rectangle"  # Default assumption
        
        shape_type = shape.shape_type
        
        # Oval/Ellipse shapes (includes circles)
        if shape_type == MSO_SHAPE_TYPE.OVAL:
            return "oval"
        
        # Placeholder, AutoShape, TextBox usually rectangular
        if shape_type in (MSO_SHAPE_TYPE.PLACEHOLDER, 
                          MSO_SHAPE_TYPE.AUTO_SHAPE,
                          MSO_SHAPE_TYPE.TEXT_BOX):
            return "rectangle"
        
        # Default to rectangle for unknown types
        return "rectangle"
    
    except Exception:
        # Fallback for any errors
        return "rectangle"
 
def get_placeholder_metadata(slide):
    """
    Extracts comprehensive metadata from all placeholders including template intelligence.
    
    Extracts:
    - Canonical identity (idx, type, name)
    - Geometry (position, size)
    - Template text and formatting
    - Font properties (size, name, style)
    - Paragraph properties (alignment, spacing)
    - Text frame properties (word wrap, margins)
    
    Returns:
        List of dicts with complete placeholder metadata
    """
    placeholders = []
    for shape in slide.placeholders:
        try:
            # Detect shape type (oval, rectangle, other)
            shape_type = map_shape_type(shape)
            
            # Start with basic metadata
            metadata = {
                # Canonical identity
                "idx": shape.placeholder_format.idx if hasattr(shape, 'placeholder_format') else None,
                "name": shape.name if hasattr(shape, 'name') else "Unnamed",
                "type": str(shape.placeholder_format.type) if hasattr(shape, 'placeholder_format') else "Unknown",
                "type_id": int(shape.placeholder_format.type) if hasattr(shape, 'placeholder_format') else None,
                
                # Geometry
                "left": shape.left if hasattr(shape, 'left') else 0,
                "top": shape.top if hasattr(shape, 'top') else 0,
                "width": shape.width if hasattr(shape, 'width') else None,
                "height": shape.height if hasattr(shape, 'height') else None,
                "shape_type": shape_type,
                
                # Template intelligence (will be populated below)
                "template_text": None,
                "default_font_size_pt": None,
                "default_font_name": None,
                "is_bold": None,
                "is_italic": None,
                "default_alignment": None,
                "line_spacing": None,
                "indent_level": None,
                "word_wrap": None,
                "auto_size": None,
                "margin_left": None,
                "margin_right": None,
                "margin_top": None,
                "margin_bottom": None
            }
            
            # Extract template intelligence from text_frame
            if hasattr(shape, 'text_frame'):
                tf = shape.text_frame
                
                # Template text (instructional content)
                if tf.text:
                    metadata["template_text"] = tf.text.strip()
                
                # Text frame properties
                metadata["word_wrap"] = tf.word_wrap if hasattr(tf, 'word_wrap') else True
                metadata["auto_size"] = str(tf.auto_size) if hasattr(tf, 'auto_size') else None
                
                # Margins (in EMU)
                if hasattr(tf, 'margin_left'):
                    metadata["margin_left"] = tf.margin_left
                if hasattr(tf, 'margin_right'):
                    metadata["margin_right"] = tf.margin_right
                if hasattr(tf, 'margin_top'):
                    metadata["margin_top"] = tf.margin_top
                if hasattr(tf, 'margin_bottom'):
                    metadata["margin_bottom"] = tf.margin_bottom
                
                # Extract font properties from first run (if exists)
                if tf.paragraphs and len(tf.paragraphs) > 0:
                    first_para = tf.paragraphs[0]
                    
                    # Paragraph properties
                    if first_para.alignment:
                        metadata["default_alignment"] = str(first_para.alignment).replace("PP_ALIGN.", "")
                    else:
                        metadata["default_alignment"] = "LEFT"
                    
                    metadata["line_spacing"] = first_para.line_spacing if hasattr(first_para, 'line_spacing') else None
                    metadata["indent_level"] = first_para.level if hasattr(first_para, 'level') else 0
                    
                    # Font properties from first run
                    if first_para.runs and len(first_para.runs) > 0:
                        first_run = first_para.runs[0]
                        font = first_run.font
                        
                        # Font size (convert to points)
                        if font.size:
                            metadata["default_font_size_pt"] = font.size.pt
                        
                        # Font name
                        metadata["default_font_name"] = font.name if font.name else None
                        
                        # Font style
                        metadata["is_bold"] = font.bold
                        metadata["is_italic"] = font.italic
            
            placeholders.append(metadata)
            
        except Exception as e:
            # Skip shapes that cause errors but log for debugging
            print(f"⚠️ Warning: Skipping problematic shape - {e}")
            continue
    
    return placeholders
 
def calculate_area(width, height):
    """
    Converts PowerPoint EMU units to a readable Area Score.
    Returns 0 if width or height is None.
    """
    if width is None or height is None:
        return 0
    return (width * height) / (914400 ** 2)

def format_text_in_placeholder(shape, content, geometry_meta):
    """
    Deterministic text formatter with geometry-based font sizing.
    Removes markdown artifacts and applies proper PowerPoint formatting.
    
    Args:
        shape: PowerPoint shape object with text_frame
        content: Raw text (may contain markdown like **bold**)
        geometry_meta: Dict with 'area_ratio', 'role_hint', 'norm_width', 'norm_height'
    
    Modifies shape in place (no return value).
    """
    if not hasattr(shape, 'text_frame'):
        return
    
    # Extract metadata
    role_hint = geometry_meta.get('role_hint', 'content')
    area_ratio = geometry_meta.get('area_ratio', 0.1)
    
    # Calculate appropriate font size based on role and geometry
    if role_hint == 'title':
        font_size = 32 if area_ratio > 0.05 else 28
    elif role_hint == 'body':
        font_size = 22 if area_ratio > 0.3 else 18
    elif role_hint == 'content':
        font_size = 16 if area_ratio > 0.15 else 14
    elif role_hint == 'footer':
        font_size = 10
    else:
        font_size = 16  # Default
    
    # Parse and clean markdown
    # Find **bold** patterns
    bold_pattern = re.compile(r'\*\*(.*?)\*\*')
    bold_segments = []
    last_end = 0
    
    for match in bold_pattern.finditer(content):
        # Add text before bold
        if match.start() > last_end:
            bold_segments.append((content[last_end:match.start()], False))
        # Add bold text
        bold_segments.append((match.group(1), True))
        last_end = match.end()
    
    # Add remaining text
    if last_end < len(content):
        bold_segments.append((content[last_end:], False))
    
    # If no markdown found, treat entire content as plain
    if not bold_segments:
        bold_segments = [(content, False)]
    
    # Clear existing text and rebuild with formatting
    text_frame = shape.text_frame
    text_frame.clear()
    
    # Get or create paragraph
    if len(text_frame.paragraphs) == 0:
        p = text_frame.paragraphs[0] if text_frame.paragraphs else text_frame.add_paragraph()
    else:
        p = text_frame.paragraphs[0]
    
    # Apply alignment based on role
    if role_hint in ['title', 'header']:
        p.alignment = PP_ALIGN.CENTER
    else:
        p.alignment = PP_ALIGN.LEFT
    
    # Add runs with formatting
    for text_segment, is_bold in bold_segments:
        run = p.add_run()
        run.text = text_segment
        run.font.size = Pt(font_size)
        if is_bold:
            run.font.bold = True
    
    # Set text frame properties for better fitting
    text_frame.word_wrap = True
    text_frame.auto_size = None  # Don't auto-resize