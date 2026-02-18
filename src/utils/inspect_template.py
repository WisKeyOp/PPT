"""
Template Inspection Utility

This script extracts and displays all placeholder metadata from a PowerPoint template.
Useful for:
- Verifying template structure before indexing
- Debugging placeholder identity issues
- Understanding template designer's intent

Usage:
    python -m src.utils.inspect_template data/templates/template2.pptx
"""
from pptx import Presentation
import sys
import os


def inspect_template(template_path):
    """
    Extract and print comprehensive placeholder metadata from a PowerPoint template.
    
    Args:
        template_path: Path to the .pptx template file
    """
    if not os.path.exists(template_path):
        print(f"❌ Error: Template not found: {template_path}")
        return
    
    try:
        prs = Presentation(template_path)
    except Exception as e:
        print(f"❌ Error loading template: {e}")
        return
    
    print(f"\n{'='*80}")
    print(f"TEMPLATE INSPECTION: {os.path.basename(template_path)}")
    print(f"{'='*80}\n")
    print(f"Slide dimensions: {prs.slide_width / 914400:.2f}\" x {prs.slide_height / 914400:.2f}\"")
    print(f"Total layouts: {len(prs.slide_layouts)}\n")
    
    for layout_idx, layout in enumerate(prs.slide_layouts):
        print(f"{'─'*80}")
        print(f"LAYOUT {layout_idx}: {layout.name}")
        print(f"{'─'*80}")
        
        if not layout.placeholders:
            print("  (No placeholders in this layout)\n")
            continue
        
        for shape in layout.placeholders:
            try:
                # Extract canonical identity
                placeholder_idx = shape.placeholder_format.idx
                placeholder_type_id = int(shape.placeholder_format.type)
                placeholder_type_name = str(shape.placeholder_format.type)
                placeholder_name = shape.name
                
                # Build stable key
                stable_key = f"{layout_idx}:{placeholder_type_name}:{placeholder_idx}"
                
                # Extract geometry
                left_in = shape.left / 914400 if hasattr(shape, 'left') else 0
                top_in = shape.top / 914400 if hasattr(shape, 'top') else 0
                width_in = shape.width / 914400 if hasattr(shape, 'width') else 0
                height_in = shape.height / 914400 if hasattr(shape, 'height') else 0
                
                # Extract text and formatting
                template_text = ""
                font_info = ""
                alignment_info = ""
                
                if hasattr(shape, 'text_frame'):
                    tf = shape.text_frame
                    
                    # Template text
                    if tf.text:
                        template_text = tf.text.strip()
                        # Truncate if too long
                        if len(template_text) > 60:
                            template_text = template_text[:60] + "..."
                    
                    # Font properties (from first run)
                    if tf.paragraphs and tf.paragraphs[0].runs:
                        first_run = tf.paragraphs[0].runs[0]
                        font = first_run.font
                        
                        font_size_pt = font.size.pt if font.size else "N/A"
                        font_name = font.name or "Default"
                        font_bold = "B" if font.bold else ""
                        font_italic = "I" if font.italic else ""
                        font_modifiers = f"{font_bold}{font_italic}".strip()
                        
                        if font_modifiers:
                            font_info = f"{font_size_pt}pt {font_name} [{font_modifiers}]"
                        else:
                            font_info = f"{font_size_pt}pt {font_name}"
                    
                    # Paragraph alignment
                    if tf.paragraphs:
                        first_para = tf.paragraphs[0]
                        if first_para.alignment:
                            alignment_info = str(first_para.alignment).replace("PP_ALIGN.", "")
                        else:
                            alignment_info = "LEFT (default)"
                
                # Print placeholder details
                print(f"\n  Placeholder {placeholder_idx}: {placeholder_name}")
                print(f"    Stable Key: [{stable_key}]")
                print(f"    Type: {placeholder_type_name} (ID: {placeholder_type_id})")
                print(f"    Geometry: {left_in:.2f}\", {top_in:.2f}\" | {width_in:.2f}\" × {height_in:.2f}\"")
                
                if template_text:
                    print(f"    Template Text: \"{template_text}\"")
                if font_info:
                    print(f"    Font: {font_info}")
                if alignment_info:
                    print(f"    Alignment: {alignment_info}")
                
            except Exception as e:
                print(f"\n  ⚠️ Placeholder {shape.name}: Error extracting metadata - {e}")
        
        print()  # Blank line after each layout
    
    print(f"{'='*80}")
    print(f"Inspection complete: {len(prs.slide_layouts)} layouts analyzed")
    print(f"{'='*80}\n")


def main():
    """Command-line entry point."""
    if len(sys.argv) < 2:
        print("Usage: python -m src.utils.inspect_template <path_to_template.pptx>")
        print("\nExample:")
        print("  python -m src.utils.inspect_template data/templates/template2.pptx")
        sys.exit(1)
    
    template_path = sys.argv[1]
    inspect_template(template_path)


if __name__ == "__main__":
    main()
